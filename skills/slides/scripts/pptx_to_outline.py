#!/usr/bin/env python3
"""
pptx_to_outline.py — turn an existing PowerPoint into a structured outline.

Extracts each slide's title, body bullets, table text, and speaker notes so the
/slides skill can map an existing deck onto its own formats and themes (rather
than re-typing it). It does NOT try to preserve PowerPoint styling — the point is
to lift the CONTENT, then rebuild it on-brand with this skill.

Usage:
    python3 pptx_to_outline.py --pptx deck.pptx              # markdown to stdout
    python3 pptx_to_outline.py --pptx deck.pptx --format json
    python3 pptx_to_outline.py --pptx deck.pptx --out outline.md

Requires python-pptx:
    pip install python-pptx
"""
import argparse, json, sys, pathlib


def shape_lines(shape):
    """Return non-empty text lines from a shape (text frame or table)."""
    lines = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
                lines.append(t)
    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
    return lines


def extract(pptx_path):
    from pptx import Presentation  # imported here so --help works without the dep
    prs = Presentation(str(pptx_path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()
        bullets = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            bullets.extend(shape_lines(shape))
        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        slides.append({"index": i, "title": title or f"Slide {i}",
                       "bullets": bullets, "notes": notes})
    return slides


def to_markdown(name, slides):
    out = [f"# {name}", f"{len(slides)} slides", ""]
    for s in slides:
        out.append(f"## Slide {s['index']} — {s['title']}")
        for b in s["bullets"]:
            out.append(f"- {b}")
        if s["notes"]:
            out.append(f"Notes: {s['notes']}")
        out.append("")
    return "\n".join(out)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True, help="path to the .pptx file")
    ap.add_argument("--format", choices=["md", "json"], default="md")
    ap.add_argument("--out", default=None, help="write to a file instead of stdout")
    args = ap.parse_args()

    pptx_path = pathlib.Path(args.pptx).expanduser().resolve()
    if not pptx_path.exists():
        sys.exit(f"ERROR: file not found: {pptx_path}")

    try:
        slides = extract(pptx_path)
    except ImportError:
        sys.exit("ERROR: python-pptx not installed. Run:\n  pip install python-pptx")

    payload = (json.dumps(slides, indent=2, ensure_ascii=False)
               if args.format == "json"
               else to_markdown(pptx_path.stem, slides))

    if args.out:
        pathlib.Path(args.out).expanduser().write_text(payload, encoding="utf-8")
        print(args.out)
    else:
        print(payload)


if __name__ == "__main__":
    main()
